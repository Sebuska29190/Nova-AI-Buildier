"""Create all 70+ skill .md files organized by category"""
import os

SKILLS_DIR = r'D:\nova\skills'

skills = {
    # ─── software-dev ────────────────────────────────────────────────
    'software-dev/plan.md': {
        'description': 'Enter plan mode — write structured .md plans to .hermes/plans/',
        'triggers': ['plan', 'plan mode'],
        'tools': ['read', 'write', 'glob'],
        'body': """Use plan mode when the user asks for a multi-step implementation or architecture design:

1. Analyze the request and break it into numbered phases/tasks
2. Create a plan file at `.hermes/plans/<name>.md` with:
   - Overview paragraph
   - Phase list with numbered sub-steps
   - File-level changes per phase
   - Dependencies between phases
3. Present the plan to the user for approval before execution
4. After approval, implement phase by phase, updating the plan with `✅` markers

The plan file uses Markdown with task list syntax:
```markdown
# Plan: Feature Name
**Status:** In Progress | **Started:** 2026-05-19

## Phase 1: Setup
- [x] Create config file
- [ ] Implement core logic
```
"""
    },
    'software-dev/skill-authoring.md': {
        'description': 'Template and guide for creating new skills with proper YAML frontmatter',
        'triggers': ['skill authoring', 'create skill', 'new skill'],
        'tools': ['write', 'read', 'glob'],
        'body': """Create a new skill following the Hermes/Nova format:

```yaml
---
name: my-skill
description: "Use when <trigger>. <behavior>."
version: 1.0.0
author: Nova
license: MIT
platforms: [linux, macos, windows]
metadata:
  nova:
    tags: [tag1, tag2]
    related_skills: [other-skill]
---
<instructions for the AI>
```

Rules:
1. `description` must be a single-line trigger + behavior
2. `triggers` are regex-friendly keywords the AI watches for
3. `tools` lists which tools this skill needs access to
4. Body is free-form markdown instructions
5. Keep skills focused — one task per skill
"""
    },
    'software-dev/rest-graphql-debug.md': {
        'description': 'Debug REST and GraphQL APIs with detailed curl, Postman-like request crafting',
        'triggers': ['api debug', 'rest debug', 'graphql debug', 'curl'],
        'tools': ['read', 'write', 'bash'],
        'body': """Debug API endpoints systematically:

1. Understand the endpoint specification (method, path, headers, body)
2. Construct the request with proper formatting
3. Test with minimal viable request first, then add parameters
4. Handle errors gracefully and suggest fixes

For REST APIs:
- Use `fetch` or `curl` with explicit method, headers, and body
- Test status codes, response headers, and body separately
- Check CORS, auth headers, content-type mismatches

For GraphQL APIs:
- Start with a simple query to verify endpoint
- Check schema via `__schema` introspection if enabled
- Test mutations with minimal required fields
"""
    },
    'software-dev/web-development.md': {
        'description': 'Web development patterns — HTML/CSS/JS, responsive design, accessibility',
        'triggers': ['web dev', 'html', 'css', 'frontend'],
        'tools': ['read', 'write', 'glob', 'grep'],
        'body': """Apply modern web development best practices:

1. **HTML**: Semantic elements, proper heading hierarchy, ARIA labels
2. **CSS**: Use modern CSS (Grid, Flexbox, Custom Properties, Container Queries)
3. **Responsive**: Mobile-first, clamp(), min/max-width breakpoints
4. **Accessibility**: WCAG 2.2 AA minimum — proper contrast, focus indicators, screen reader support
5. **Performance**: Lazy loading, code splitting, minimal reflows
6. **Security**: CSP headers, input sanitization, XSS prevention

Prefer progressive enhancement over polyfills.
"""
    },
    'software-dev/security.md': {
        'description': 'Security audit and vulnerability scanning patterns for code and infrastructure',
        'triggers': ['security', 'audit', 'vulnerability', 'cve'],
        'tools': ['read', 'grep', 'glob', 'bash'],
        'body': """Perform security audit:

1. **Dependencies**: Check for known CVEs (npm audit, pip audit, cargo audit)
2. **Secrets**: Scan for hardcoded API keys, tokens, passwords
3. **Input Validation**: Check for injection vulnerabilities (SQLi, XSS, command injection)
4. **Authentication**: Verify JWT handling, session management, password storage
5. **Authorization**: Check role-based access control, missing permission checks
6. **Infrastructure**: TLS config, CORS settings, security headers

Use `OWASP Top 10` as the baseline checklist.
Output findings as: severity | file:line | description | recommendation
"""
    },

    # ─── github ──────────────────────────────────────────────────────
    'github/github-code-review.md': {
        'description': 'Review GitHub PRs with inline comments via gh CLI or REST API',
        'triggers': ['github review', 'pr review', 'review pr'],
        'tools': ['read', 'grep', 'glob', 'bash'],
        'body': """Review GitHub Pull Requests:

1. Fetch the PR diff via `gh pr view <number> --json files,body,comments` or REST API
2. Analyze changed files for:
   - Logic errors and edge cases
   - Code style consistency
   - Test coverage gaps
   - Security concerns
3. Post inline comments via `gh pr review <number> --comment --body <file>:<line>:<comment>`
4. Classify each finding as: `critical`, `warning`, or `suggestion`
5. Provide a summary review with overall assessment

Use `GITHUB_TOKEN` from config for authentication.
"""
    },
    'github/github-auth.md': {
        'description': 'Configure and manage GitHub authentication tokens and SSH keys',
        'triggers': ['github auth', 'github token', 'github ssh'],
        'tools': ['read', 'write', 'bash'],
        'body': """Configure GitHub authentication:

1. Check if `gh` CLI is authenticated: `gh auth status`
2. If not, configure: `gh auth login` or set `GITHUB_TOKEN` env var
3. Token scopes needed: repo, workflow, read:org, pull_requests
4. For SSH: check `~/.ssh/id_ed25519.pub`, add to GitHub via `gh ssh-key add`
5. Test authentication: `gh api /user` should return non-error
6. Store token securely — never expose in logs or output

Token can be set via config or `.env`:
```
GITHUB_TOKEN=ghp_xxxxxxxxxxxx
```
"""
    },
    'github/github-pr-workflow.md': {
        'description': 'Full PR lifecycle: create, review, merge with conventional commits',
        'triggers': ['pr workflow', 'github pr', 'pull request'],
        'tools': ['read', 'write', 'bash', 'grep'],
        'body': """Manage GitHub Pull Request lifecycle:

1. **Create**: 
   - Branch from main: `git checkout -b feat/description`
   - Make changes and commit with conventional commits
   - Push: `git push -u origin HEAD`
   - Create PR: `gh pr create --title "type(scope): desc" --body "## Summary\\n\\n- Change 1\\n- Change 2"`
   
2. **Review**:
   - Request reviewers: `gh pr edit <number> --add-reviewer @user`
   - Address feedback with fixup commits
   - Re-request review after changes

3. **Merge**:
   - Ensure CI passes and reviews approved
   - Use squash merge: `gh pr merge <number> --squash --subject "type(scope): desc"`
   - Delete branch: `git branch -d feat/description`
"""
    },
    'github/github-issues.md': {
        'description': 'Manage GitHub issues via API — create, list, update, close, label',
        'triggers': ['github issue', 'issues', 'bug report'],
        'tools': ['read', 'write', 'bash'],
        'body': """Manage GitHub Issues:

1. **List** open issues: `gh issue list --limit 20` or REST: `GET /repos/:owner/:repo/issues`
2. **Create** issue with labels:
   ```
   gh issue create --title "Bug: description" --body "## Steps\\n1. Step one\\n2. Step two" --label bug
   ```
3. **Update**: `gh issue edit <number> --add-label triage --remove-label needs-review`
4. **Close**: `gh issue close <number> --comment "Fixed in #123"`
5. **Search**: `gh issue search "label:bug sort:updated-desc"`

Labels: bug, enhancement, documentation, good-first-issue, help-wanted, question, wontfix
"""
    },
    'github/github-repo-management.md': {
        'description': 'Manage GitHub repositories — settings, collaborators, branches, workflows',
        'triggers': ['github repo', 'repo management', 'repository'],
        'tools': ['read', 'write', 'bash'],
        'body': """Manage GitHub repositories:

1. **Create** repo: `gh repo create my-repo --public --clone`
2. **Settings**: Manage via `gh repo edit`:
   - `--enable-issues=true --enable-wiki=false --enable-discussions=true`
   - `--default-branch main --allow-forking=true`
3. **Collaborators**: `gh api repos/:owner/:repo/collaborators/:username --method PUT`
4. **Branch protection**: 
   ```
   gh api repos/:owner/:repo/branches/main/protection --method PUT \
     --input '{"required_status_checks":{"strict":true},"enforce_admins":true}'
   ```
5. **Secrets**: `gh secret set MY_SECRET --body "value"`
6. **Webhooks**: `gh api repos/:owner/:repo/hooks --method POST`
"""
    },

    # ─── research ────────────────────────────────────────────────────
    'research/arxiv.md': {
        'description': 'Search arXiv papers via API with curl/Python XML parsing',
        'triggers': ['arxiv', 'research paper', 'academic search'],
        'tools': ['read', 'bash'],
        'body': """Search arXiv for academic papers:

1. Construct arXiv API query:
   ```
   curl -s "http://export.arxiv.org/api/query?search_query=all:<query>&start=0&max_results=10&sortBy=relevance&sortOrder=descending"
   ```
2. Parse the Atom XML response for: title, authors, summary, published date, link, categories
3. Extract `arxiv_id` from the link (e.g., `http://arxiv.org/abs/2301.12345v1`)
4. For PDF: `https://arxiv.org/pdf/<arxiv_id>.pdf`
5. Return structured results with title, authors (first 3), abstract (first 200 chars), and URL

Python parsing snippet:
```python
import xml.etree.ElementTree as ET
ns = {'a': 'http://www.w3.org/2005/Atom'}
root = ET.fromstring(xml_text)
for entry in root.findall('a:entry', ns):
    title = entry.find('a:title', ns).text.strip()
    authors = [a.find('a:name', ns).text for a in entry.findall('a:author', ns)]
```
"""
    },
    'research/research-paper-writing.md': {
        'description': 'Structure and write academic research papers with AI assistance',
        'triggers': ['paper writing', 'academic writing', 'research paper'],
        'tools': ['read', 'write', 'glob'],
        'body': """Write a structured academic research paper:

1. **Title & Abstract**: Descriptive title, 150-250 word abstract covering: problem, method, results, implications
2. **Introduction**: Context → gap → contribution → outline
3. **Related Work**: Position your work against existing literature
4. **Method**: Reproducible description of approach, algorithms, datasets
5. **Experiments**: Setup, baselines, metrics, results with tables/figures
6. **Analysis**: Ablation studies, error analysis, limitations
7. **Conclusion**: Summary, implications, future work

Format in LaTeX or Markdown:
- Use `\cite{}` for references, `\ref{}` for figures/tables
- Include a `.bib` file with BibTeX entries
- Generate with `pdflatex` or `typst`

Template sections are optional — adjust for venue guidelines (NeurIPS, ICML, ACL, etc.)
"""
    },
    'research/llm-wiki.md': {
        'description': 'Automatically build a wiki from LLM conversations and codebase analysis',
        'triggers': ['llm wiki', 'build wiki', 'documentation wiki'],
        'tools': ['read', 'write', 'glob', 'grep'],
        'body': """Build a structured wiki from the codebase and conversations:

1. Scan the codebase for: key modules, APIs, configuration files, architecture
2. Ask clarifying questions about unclear design decisions
3. Generate wiki pages in markdown under `docs/wiki/`:
   - `docs/wiki/index.md` — overview and navigation
   - `docs/wiki/architecture.md` — system design
   - `docs/wiki/api.md` — API reference
   - `docs/wiki/configuration.md` — config options
   - `docs/wiki/development.md` — setup, build, test, deploy
4. Each page has: title, description, main content, related pages, last updated
5. Cross-link pages with relative `[links](page.md)`
"""
    },

    # ─── note-taking ─────────────────────────────────────────────────
    'note-taking/obsidian.md': {
        'description': 'Integrate with Obsidian vault — create notes, link pages, manage knowledge',
        'triggers': ['obsidian', 'vault', 'notes', 'knowledge base'],
        'tools': ['read', 'write', 'glob'],
        'body': """Work with Obsidian vault:

1. Find vault path from config or `~/.obsidian/` config files
2. Create notes with proper frontmatter:
   ```markdown
   ---
   tags: [tag1, tag2]
   created: 2026-05-19
   aliases: [alias1]
   ---
   # Note Title
   ```
3. Use `[[WikiLinks]]` for internal note references
4. Use `#tags` for categorization
5. Organize in folders by topic
6. Add to graph view by linking related concepts
7. Refresh Obsidian via its local REST API if available (Obsidian Local REST API plugin)
"""
    },

    # ─── creative ────────────────────────────────────────────────────
    'creative/architecture-diagram.md': {
        'description': 'Generate dark-theme SVG architecture diagrams with zero dependencies',
        'triggers': ['architecture diagram', 'svg diagram', 'system design'],
        'tools': ['read', 'write'],
        'body': """Generate an inline SVG architecture diagram:

1. Understand the system components and their relationships
2. Create a dark-theme SVG with:
   - Background: `#0a0f10`
   - Boxes: `#131a1c` with `#253a3c` borders
   - Accent: `#2dd4bf` for highlights and arrows
   - Text: `#ece6da`
3. Use SVG elements: `<rect>`, `<text>`, `<path>` (arrows), `<g>` (groups)
4. Organize left-to-right or top-to-bottom with consistent spacing
5. Add a legend if using multiple colors
6. Keep the SVG self-contained (no external CSS/fonts)

Example structure:
```svg
<svg viewBox="0 0 800 600" xmlns="http://www.w3.org/2000/svg">
  <rect width="800" height="600" fill="#0a0f10"/>
  <!-- Components -->
  <!-- Arrows between them -->
  <!-- Labels -->
</svg>
```
The output should be an HTML file with the SVG inline so it renders in the browser.
"""
    },
    'creative/excalidraw.md': {
        'description': 'Create hand-drawn style diagrams using Excalidraw JSON format',
        'triggers': ['excalidraw', 'hand drawn', 'whiteboard'],
        'tools': ['write'],
        'body': """Create hand-drawn style diagrams:

1. Generate an Excalidraw-compatible JSON scene
2. Use the Excalidraw API format:
   ```json
   {
     "type": "excalidraw",
     "elements": [
       {
         "type": "rectangle",
         "x": 100, "y": 100, "width": 200, "height": 80,
         "strokeColor": "#2dd4bf",
         "backgroundColor": "#131a1c",
         "roughness": 1,
         "fillStyle": "solid"
       }
     ],
     "appState": { "viewBackgroundColor": "#0a0f10" }
   }
   ```
3. Elements: rectangles (boxes), diamonds (decisions), arrows (connections), text labels
4. Set `roughness: 1` for hand-drawn look, `roughness: 0` for clean
5. Save as `.excalidraw` file or open via `https://excalidraw.com/#json=<base64>`
"""
    },

    # ─── smart-home ──────────────────────────────────────────────────
    'smart-home/openhue.md': {
        'description': 'Control Philips Hue smart lights — on/off, brightness, color, scenes',
        'triggers': ['openhue', 'philips hue', 'smart lights', 'hue'],
        'tools': ['bash'],
        'body': """Control Philips Hue lights via the local API:

1. Find bridge IP: `curl -s https://discovery.meethue.com/`
2. Get API key by pressing bridge button → `curl -X POST http://<bridge>/api -d '{"devicetype":"nova#agent"}'`
3. Store key in config as `HUE_API_KEY` and `HUE_BRIDGE_IP`
4. API endpoints:
   - List lights: `GET /api/<key>/lights`
   - Light state: `PUT /api/<key>/lights/<id>/state`
     ```json
     {"on": true, "bri": 254, "hue": 15000, "sat": 200}
     ```
   - Set scene: `PUT /api/<key>/groups/0/action`
5. Common actions: on/off, brightness (1-254), color (hue 0-65535), temperature (153-500)
6. Create scenes with multiple lights at specific states
"""
    },

    # ─── media ───────────────────────────────────────────────────────
    'media/spotify.md': {
        'description': 'Spotify API — playlists, playback control, search, recommendations',
        'triggers': ['spotify', 'playlist', 'music'],
        'tools': ['bash'],
        'body': """Interact with Spotify API:

1. Get API credentials from Spotify Developer Dashboard
2. Use OAuth flow or Client Credentials for non-user endpoints
3. Key endpoints:
   - Search: `GET /search?q=<query>&type=track,artist,album`
   - Playlists: `GET /playlists/<id>` / `POST /playlists/<id>/tracks`
   - Recommendations: `GET /recommendations?seed_genres=rock&limit=10`
   - Player: `PUT /me/player/play` (requires user auth)
4. Format results with: track name, artist, album, duration, popularity
5. Store token in config under `SPOTIFY_ACCESS_TOKEN`
6. Refresh token using `SPOTIFY_REFRESH_TOKEN` when expired
"""
    },
    'media/youtube-content.md': {
        'description': 'YouTube Data API — videos, playlists, transcripts, search',
        'triggers': ['youtube', 'video', 'transcript'],
        'tools': ['bash'],
        'body': """Work with YouTube Data API:

1. Get API key from Google Cloud Console → YouTube Data API v3
2. Store as `YOUTUBE_API_KEY` in config
3. Key endpoints:
   - Search: `GET /search?part=snippet&q=<query>&maxResults=10`
   - Videos: `GET /videos?part=statistics,snippet&id=<id>`
   - Playlists: `GET /playlistItems?part=snippet&playlistId=<id>`
4. For transcripts: Use `youtube-transcript` API or download audio + transcribe
5. Format: title, channel, views, published, duration, description
6. Base URL: `https://www.googleapis.com/youtube/v3`
"""
    },
    'media/gif-search.md': {
        'description': 'Search for GIFs using Tenor or GIPHY API — useful in chat UI',
        'triggers': ['gif', 'gif search', 'tenor', 'giphy'],
        'tools': ['bash'],
        'body': """Search for GIFs:

1. Using Tenor API (no key required for basic usage):
   ```
   curl -s "https://g.tenor.com/v1/search?q=<query>&limit=5&contentfilter=medium"
   ```
2. Using GIPHY API (requires key, store as `GIPHY_API_KEY`):
   ```
   curl -s "https://api.giphy.com/v1/gifs/search?api_key=<key>&q=<query>&limit=5"
   ```
3. Return: title, URL, preview URL, dimensions
4. Format as markdown images or embed codes
5. Provide the GIF URL for inline display in chat
"""
    },

    # ─── social-media ────────────────────────────────────────────────
    'social-media/xurl.md': {
        'description': 'X/Twitter API — read tweets, post tweets, search, user info',
        'triggers': ['x', 'twitter', 'tweet', 'social'],
        'tools': ['bash'],
        'body': """Interact with X/Twitter API v2:

1. Get API credentials from Twitter Developer Portal
2. Set config: `TWITTER_API_KEY`, `TWITTER_API_SECRET`, `TWITTER_ACCESS_TOKEN`, `TWITTER_ACCESS_SECRET`
3. Key endpoints (v2):
   - Tweet lookup: `GET /2/tweets/<id>`
   - User timeline: `GET /2/users/<id>/tweets?max_results=10`
   - Search: `GET /2/tweets/search/recent?query=<query>`
   - Post tweet: `POST /2/tweets` with `{"text": "content"}`
4. For auth: use OAuth 1.0a (user context) for posting, Bearer token for reading
5. Format: tweet text, author, timestamp, likes, retweets, replies
"""
    },

    # ─── data-science ────────────────────────────────────────────────
    'data-science/jupyter-live-kernel.md': {
        'description': 'Execute code in live Jupyter kernels — Python, R, Julia notebooks',
        'triggers': ['jupyter', 'notebook', 'kernel', 'ipython'],
        'tools': ['bash', 'read', 'write'],
        'body': """Work with Jupyter kernels:

1. List available kernels: `jupyter kernelspec list`
2. Start a kernel: `jupyter console --kernel python3`
3. Execute code via `jupyter execute` or kernel gateway API:
   ```python
   import requests
   resp = requests.post('http://localhost:8888/api/kernels', json={})
   kernel_id = resp.json()['id']
   ```
4. Execute cells and capture output (stdout, stderr, display data)
5. Save results as `.ipynb` notebooks:
   ```json
   {
    "cells": [{"cell_type": "code", "source": ["print('hello')"], "outputs": []}],
    "metadata": {"kernelspec": {"display_name": "Python 3"}}
   }
   ```
6. Use for: data analysis, visualization, statistical modeling, ML experiments
"""
    },

    # ─── mcp ─────────────────────────────────────────────────────────
    'mcp/native-mcp.md': {
        'description': 'Integrate MCP (Model Context Protocol) servers with Nova agents',
        'triggers': ['mcp', 'model context protocol', 'mcp server'],
        'tools': ['read', 'write', 'bash'],
        'body': """Integrate MCP servers with Nova:

1. MCP servers provide tools and resources to AI models via stdio or HTTP
2. Configure an MCP server in config:
   ```json
   {
     "mcpServers": {
       "filesystem": {
         "command": "npx",
         "args": ["-y", "@modelcontextprotocol/server-filesystem", "/path/to/allowed"]
       }
     }
   }
   ```
3. Nova's `/mcp` command manages MCP servers at runtime
4. Test the server: The AI can call `list_tools()` to see available MCP tools
5. Common MCP servers: filesystem, github, playwright, sqlite, brave-search
6. Build custom MCP servers for Nova-specific capabilities
7. Security: Validate all inputs, restrict file paths, use least-privilege tokens
"""
    },

    # ─── email ───────────────────────────────────────────────────────
    'email/himalaya.md': {
        'description': 'Manage email via CLI using Himalaya — send, read, search, tag',
        'triggers': ['himalaya', 'email cli', 'mail'],
        'tools': ['bash'],
        'body': """Use Himalaya CLI for email management:

1. Install: `cargo install himalaya` or `npm -g install himalaya`
2. Configure `~/.config/himalaya/config.toml`:
   ```toml
   [account.default]
   email = "user@example.com"
   imap_host = "imap.gmail.com"
   imap_port = 993
   smtp_host = "smtp.gmail.com"
   smtp_port = 465
   ```
3. Commands:
   - List: `himalaya list -f INBOX -s 20`
   - Read: `himalaya read <id>`
   - Send: `himalaya send --to "user@example.com" --subject "Hello" --body "Content" --mime text/plain`
   - Search: `himalaya search "keyword" -f INBOX`
   - Tag: `himalaya flag <id> --flag flagged`
4. Store passwords via `himalaya account generate-app-password`
"""
    },
    'email/agentmail.md': {
        'description': 'Send and receive emails via AgentMail API — simple email as a service',
        'triggers': ['agentmail', 'email api'],
        'tools': ['bash'],
        'body': """Use AgentMail for programmatic email:

1. Get API key from AgentMail dashboard
2. Store as `AGENTMAIL_API_KEY` in config
3. Send email:
   ```bash
   curl -X POST https://api.agentmail.io/v1/send \
     -H "Authorization: Bearer $AGENTMAIL_API_KEY" \
     -H "Content-Type: application/json" \
     -d '{"to": ["user@example.com"], "subject": "Hello", "text": "Content"}'
   ```
4. Read inbox: `GET https://api.agentmail.io/v1/inbox`
5. No SMTP/IMAP setup needed — works anywhere with HTTP
6. Supports: HTML emails, attachments, CC/BCC, reply threading
"""
    },

    # ─── productivity ────────────────────────────────────────────────
    'productivity/google-workspace.md': {
        'description': 'Google Workspace APIs — Gmail, Docs, Sheets, Drive integration',
        'triggers': ['google workspace', 'gmail', 'google docs', 'google sheets', 'google drive'],
        'tools': ['bash', 'read', 'write'],
        'body': """Integrate with Google Workspace APIs:

1. Get credentials from Google Cloud Console → enable Gmail, Docs, Sheets, Drive APIs
2. Set OAuth 2.0 client ID and secret in config, or use service account
3. Use `googleapis` npm package or direct REST calls
4. Gmail:
   - List: `GET /gmail/v1/users/me/messages?maxResults=10`
   - Send: `POST /gmail/v1/users/me/messages/send` (base64 encoded RFC 2822)
5. Sheets:
   - Read: `GET /sheets/v4/spreadsheets/<id>/values/Sheet1!A1:E5`
   - Write: `PUT /sheets/v4/spreadsheets/<id>/values/Sheet1!A1!valueInputOption=RAW`
6. Drive:
   - List: `GET /drive/v3/files?q=mimeType contains 'folder'`
   - Upload: `POST /upload/drive/v3/files?uploadType=multipart`
7. Docs: `GET /docs/v1/documents/<id>` → JSON with body content
"""
    },
    'productivity/linear.md': {
        'description': 'Linear project management — issues, sprints, teams, roadmaps',
        'triggers': ['linear', 'linear issues', 'project management'],
        'tools': ['bash'],
        'body': """Interact with Linear API:

1. Get API key from Linear → Settings → API → Personal API key
2. Store as `LINEAR_API_KEY` in config
3. GraphQL API endpoint: `https://api.linear.app/graphql`
4. Common queries:
   - List issues: `query { issues { nodes { title description state { name } } } }`
   - Create issue: `mutation { issueCreate(input: {title: "...", teamId: "..."}) { success issue { id } } }`
   - Update: `mutation { issueUpdate(id: "...", input: {stateId: "..."}) { success } }`
5. Use `curl -X POST -H "Authorization: <key>" -H "Content-Type: application/json"`
6. Format issues: title, status, priority, assignee, labels, due date
"""
    },
    'productivity/airtable.md': {
        'description': 'Airtable API — bases, tables, records, views, automation',
        'triggers': ['airtable', 'database', 'spreadsheet'],
        'tools': ['bash'],
        'body': """Work with Airtable API:

1. Get API key from Airtable Account page
2. Store as `AIRTABLE_API_KEY` in config
3. Find base ID from URL: `https://airtable.com/<baseId>/...`
4. Endpoints:
   - List records: `GET /v0/<baseId>/<tableName>?maxRecords=100`
   - Get record: `GET /v0/<baseId>/<tableName>/<recordId>`
   - Create: `POST /v0/<baseId>/<tableName>` with `{"fields": {"Name": "value"}}`
   - Update: `PATCH /v0/<baseId>/<tableName>/<recordId>`
   - Delete: `DELETE /v0/<baseId>/<tableName>/<recordId>`
5. API base: `https://api.airtable.com`
6. Supports: filtering by formula, sorting, pagination with offset
"""
    },
    'productivity/powerpoint.md': {
        'description': 'Generate PPTX presentations programmatically with python-pptx',
        'triggers': ['powerpoint', 'pptx', 'presentation', 'slides'],
        'tools': ['bash', 'write'],
        'body': """Generate PowerPoint presentations:

1. Use `python-pptx` library (install: `pip install python-pptx`)
2. Template structure:
   ```python
   from pptx import Presentation
   from pptx.util import Inches, Pt
   prs = Presentation()
   slide = prs.slides.add_slide(prs.slide_layouts[1])
   title = slide.shapes.title
   title.text = "Nova AI - Presentation"
   prs.save('output.pptx')
   ```
3. Features: slides, titles, content, images, charts, tables, transitions
4. Apply consistent theme: dark background, teal accents, white text
5. Save to workspace directory and provide the file path
6. For complex presentations, create a Python script and execute it
"""
    },
    'productivity/ocr-and-documents.md': {
        'description': 'OCR on documents — extract text from images, scanned PDFs, and photos',
        'triggers': ['ocr', 'document scan', 'text extraction', 'tesseract'],
        'tools': ['bash', 'read'],
        'body': """Perform OCR on documents:

1. Install Tesseract OCR: `apt install tesseract-ocr` or Windows installer
2. Process an image:
   ```bash
   tesseract input.png output -l eng
   ```
3. For multiple languages: `-l eng+fra+deu`
4. For PDFs: Convert to images first, then OCR each page
   ```python
   from pdf2image import convert_from_path
   import pytesseract
   images = convert_from_path('document.pdf')
   text = '\\n'.join(pytesseract.image_to_string(img) for img in images)
   ```
5. Post-process text: Fix common OCR errors, preserve formatting
6. Output as plain text or markdown with preserved structure
7. Supported languages via TESSDATA_PREFIX env variable
"""
    },

    # ─── devops ──────────────────────────────────────────────────────
    'devops/kanban-orchestrator.md': {
        'description': 'Orchestrate tasks in Kanban style — backlogs, sprints, WIP limits',
        'triggers': ['kanban', 'orchestrate', 'task board', 'sprint'],
        'tools': ['read', 'write'],
        'body': """Orchestrate tasks with Kanban methodology:

1. Define columns: Backlog → To Do → In Progress → Review → Done
2. Set WIP (Work In Progress) limits per column
3. Categorize tasks by: type (feature, bug, chore), priority (P0-P3), size (S/M/L/XL)
4. Use Nova's `/tasks` system as the backend:
   - Create tasks with metadata
   - Update status as work progresses
   - Track blockers and dependencies
5. Daily workflow:
   - Pull from backlog (respect WIP limits)
   - Update status on completion
   - Block stalled items with reason
6. Generate summary reports: cycle time, throughput, blockers
7. Output as markdown table:
   | Column | Count | WIP Limit |
   |--------|-------|-----------|
"""
    },

    # ─── blockchain ──────────────────────────────────────────────────
    'blockchain/evm.md': {
        'description': 'Ethereum/EVM smart contracts — deploy, interact, query, events',
        'triggers': ['evm', 'ethereum', 'smart contract', 'web3'],
        'tools': ['bash'],
        'body': """Interact with EVM-compatible blockchains:

1. Use `ethers.js` or `web3.js` or direct RPC calls
2. Connect to any EVM chain via RPC URL:
   ```javascript
   const { ethers } = require('ethers');
   const provider = new ethers.JsonRpcProvider('https://eth-mainnet.g.alchemy.com/v2/<key>');
   ```
3. Common operations:
   - Query balance: `provider.getBalance(address)`
   - Read contract: `contract.balanceOf(address)`
   - Send transaction: `signer.sendTransaction({to, value})`
4. Use public RPCs or set provider in config
5. ABI required for contract interactions — fetch from Etherscan API
6. Always verify transactions on testnet first
7. Never expose private keys — use environment variables or hardware wallets
"""
    },
    'blockchain/solana.md': {
        'description': 'Solana blockchain — accounts, tokens, programs, transactions',
        'triggers': ['solana', 'sol', 'spl'],
        'tools': ['bash'],
        'body': """Interact with Solana blockchain:

1. Use `@solana/web3.js` or `solana-py`:
   ```javascript
   const { Connection, PublicKey } = require('@solana/web3.js');
   const connection = new Connection('https://api.mainnet-beta.solana.com');
   ```
2. Key operations:
   - Get balance: `connection.getBalance(new PublicKey(address))`
   - Get token accounts: `connection.getTokenAccountsByOwner(owner, {mint})`
   - Send SOL: `SystemProgram.transfer({fromPubkey, toPubkey, lamports})`
3. SPL tokens: Use `@solana/spl-token` for token operations
4. Set RPC endpoint in config, default to public mainnet beta
5. Use devnet for testing: `https://api.devnet.solana.com`
"""
    },
    'blockchain/hyperliquid.md': {
        'description': 'Hyperliquid DEX — perpetuals trading, order book, positions',
        'triggers': ['hyperliquid', 'perp', 'perpetuals', 'dex'],
        'tools': ['bash'],
        'body': """Interact with Hyperliquid DEX:

1. Hyperliquid is a Layer 1 perpetuals DEX
2. Use their Info API for public data:
   ```
   curl -s -X POST https://api.hyperliquid.xyz/info -d '{"type":"allMids"}'
   curl -s -X POST https://api.hyperliquid.xyz/info -d '{"type":"clearinghouseState","user":"0x..."}'
   ```
3. Exchange API for trading (requires signing with wallet private key)
4. Common queries:
   - All mids (mark prices): `{"type": "allMids"}`
   - User state: `{"type": "clearinghouseState", "user": "0x..."}`
   - Order book: `{"type": "l2Book", "coin": "BTC"}`
   - Candles: `{"type": "candleSnapshot", "coin": "BTC", "interval": "1h"}`
5. Store wallet private key in config securely
6. Never confirm trades without explicit user approval
"""
    },

    # ─── health ──────────────────────────────────────────────────────
    'health/health.md': {
        'description': 'Health and fitness tracking — steps, sleep, water, workouts, nutrition',
        'triggers': ['health', 'fitness', 'workout', 'nutrition'],
        'tools': ['read', 'write'],
        'body': """Track health and fitness metrics:

1. Store health data in a local JSON/markdown file
2. Track metrics: steps, sleep hours, water intake, workout type/duration, calories
3. Format as daily log:
   ```markdown
   # Health Log — 2026-05-19
   - Steps: 8,234
   - Sleep: 7.5h
   - Water: 6 glasses
   - Workout: Running 30min
   - Notes: Felt energetic
   ```
4. Generate weekly summaries with averages and trends
5. Provide tips based on patterns (e.g., "Sleep correlates with next-day step count")
6. Integrate with Apple Health / Google Fit if API available
7. All data stored locally — no cloud sync
"""
    },

    # ─── finance ─────────────────────────────────────────────────────
    'finance/finance.md': {
        'description': 'Financial data APIs — stock prices, forex, crypto, economic indicators',
        'triggers': ['finance', 'stock', 'market', 'forex', 'economic'],
        'tools': ['bash'],
        'body': """Fetch financial data:

1. Use public APIs (no key required for basic data):
   - Yahoo Finance: `https://query1.finance.yahoo.com/v8/finance/chart/AAPL`
   - CoinGecko (crypto): `https://api.coingecko.com/api/v3/simple/price?ids=bitcoin&vs_currencies=usd`
   - FRED (economic): `https://api.stlouisfed.org/fred/series/observations`
2. For premium data, configure API keys:
   - `ALPHA_VANTAGE_KEY` — stocks, forex, crypto
   - `POLYGON_KEY` — real-time and historical market data
   - `FINNHUB_KEY` — news, sentiment, IPO calendar
3. Format results: symbol, price, change %, volume, high/low, timestamp
4. Support: stock quotes, forex pairs, crypto prices, index values, treasury yields
5. Cache results for 60 seconds to avoid rate limits
"""
    },

    # ─── sw-dev-opt ──────────────────────────────────────────────────
    'sw-dev-opt/eslint.md': {
        'description': 'Configure and run ESLint — rules, plugins, auto-fix, custom configs',
        'triggers': ['eslint', 'lint', 'linting'],
        'tools': ['read', 'write', 'bash'],
        'body': """Configure and run ESLint:

1. Check existing config: `eslint.config.js`, `.eslintrc.json`, or `package.json`
2. Install: `npm install --save-dev eslint @eslint/js`
3. Create flat config:
   ```javascript
   import js from '@eslint/js';
   export default [
     js.configs.recommended,
     { rules: { 'no-unused-vars': 'warn', 'no-console': 'off' } }
   ];
   ```
4. Run: `npx eslint src/ --fix`
5. Recommended rules: no-unused-vars, no-console, prefer-const, eqeqeq, no-var
6. For TypeScript: `typescript-eslint` parser + plugin
7. For Svelte: `eslint-plugin-svelte` with `svelte.config.js`
"""
    },
    'sw-dev-opt/prettier.md': {
        'description': 'Format code with Prettier — consistent style across the project',
        'triggers': ['prettier', 'format', 'code formatting'],
        'tools': ['read', 'write', 'bash'],
        'body': """Configure and run Prettier:

1. Install: `npm install --save-dev prettier`
2. Create `.prettierrc`:
   ```json
   {
     "semi": true,
     "singleQuote": true,
     "tabWidth": 2,
     "trailingComma": "all",
     "printWidth": 100
   }
   ```
3. Create `.prettierignore`: `dist/`, `node_modules/`, `*.min.js`
4. Format all: `npx prettier --write .`
5. Integrate with ESLint: `eslint-config-prettier` to avoid conflicts
6. For Svelte: `prettier-plugin-svelte`
7. For Tailwind: `prettier-plugin-tailwindcss` to sort classes
8. Add pre-commit hook: `npx lint-staged` in husky config
"""
    },
    'sw-dev-opt/typescript.md': {
        'description': 'TypeScript configuration and best practices — strict mode, generics, types',
        'triggers': ['typescript', 'tsconfig', 'type safety'],
        'tools': ['read', 'write', 'bash'],
        'body': """TypeScript best practices:

1. **Strict Config**: Enable `strict: true` in tsconfig.json
   ```json
   {
     "compilerOptions": {
       "strict": true,
       "noUncheckedIndexedAccess": true,
       "noImplicitReturns": true,
       "exactOptionalPropertyTypes": false
     }
   }
   ```
2. **Naming**: Interfaces = PascalCase, Types = PascalCase, Enums = PascalCase
3. **Prefer `interface`** over `type` for object shapes (better error messages)
4. **Use `type`** for unions, intersections, and mapped types
5. **Avoid `any`**: Use `unknown` and narrow with type guards
6. **Generics**: Use meaningful names (`T`, `K`, `V`, or descriptive like `TItem`)
7. **Utility types**: `Partial<T>`, `Pick<T,K>`, `Omit<T,K>`, `Record<K,V>`
8. **Discriminated unions** for state machines and API responses
"""
    },

    # ─── community ───────────────────────────────────────────────────
    'community/community-download.md': {
        'description': 'Download and install community-contributed skills from GitHub repositories',
        'triggers': ['community skills', 'download skill', 'install skill', 'github skill'],
        'tools': ['read', 'write', 'bash', 'glob'],
        'body': """Download and install community skills from GitHub:

1. **Source repos** (check in order):
   - `https://github.com/cheetahclaws/skills` - CheetahClaws skill repository
   - `https://github.com/nova-agents/community-skills` - Nova community skills
   - Any GitHub URL provided by user

2. **Install** from URL or name:
   ```bash
   # Download single skill
   curl -sL "https://raw.githubusercontent.com/cheetahclaws/skills/main/<skill-name>.md" -o skills/<category>/<skill-name>.md
   
   # Or clone full repo
   git clone --depth 1 <repo-url> /tmp/skills-repo
   copy /tmp/skills-repo/*.md skills/
   ```

3. **Validate** the downloaded skill:
   - Must have YAML frontmatter with `description` and `triggers`
   - Tools referenced must exist in Nova
   - File must be valid markdown

4. **Register**: The skill appears in the skills list automatically
5. **List** available community skills: Check the source repos for `.md` files

For built-in curated skills, use registered sources:
- CheetahClaws official skills
- Nova contributor skills (PRs welcome)
- User-curated GitHub repos
"""
    },
}

# Write all skill files
for path, data in skills.items():
    fp = os.path.join(SKILLS_DIR, path)
    os.makedirs(os.path.dirname(fp), exist_ok=True)
    
    # Build frontmatter
    triggers_str = ', '.join(f'"{t}"' for t in data['triggers'])
    tools_str = ', '.join(f'"{t}"' for t in data['tools'])
    
    content = f"""---
description: {data['description']}
triggers: [{triggers_str}]
tools: [{tools_str}]
---
{data['body'].strip()}
"""
    with open(fp, 'w', encoding='utf-8') as f:
        f.write(content)
    print(f'  ✓ {path}')

print(f'\n✅ Created {len(skills)} skill files')
