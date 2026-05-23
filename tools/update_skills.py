"""Update all skill .md files with proper detailed bodies"""
import os

SKILLS_DIR = r'D:\nova\skills'

# Detailed bodies for each skill
bodies = {
    # ─── software-dev ───────────────────────────────────────────
    'software-dev/plan.md': """Use plan mode when the user asks for a multi-step implementation or architecture design.

## Workflow
1. Analyze the request and break it into numbered phases/tasks
2. Use the `plan_save` tool to save the plan as a structured .md file
3. Present the plan to the user for approval before execution
4. After approval, implement phase by phase, updating the plan with checkmarks

## Plan Format
```markdown
# Plan: Feature Name
**Status:** In Progress

## Phase 1: Setup
- [x] Create config
- [ ] Implement core logic
- [ ] Write tests
```

## Output
- Save to `.cheetahclaws/plans/<name>.md`
- Use `plan_list` to see all plans
""",

    'software-dev/skill-authoring.md': """Create a new skill following the Hermes/Nova format. A skill is a `.md` file with YAML frontmatter that tells the AI when to activate and what to do.

## Template
```yaml
---
name: my-skill
description: "Use when <trigger>. <behavior>."
version: 1.0.0
author: Nova
license: MIT
platforms: [linux, macos, windows]
metadata:
  nova:
    tags: [category]
    related_skills: [other-skill]
---
<instructions for the AI>
```

## Rules
1. description must be a one-liner trigger + behavior
2. triggers are keywords the AI watches for in user queries
3. tools lists which Nova tools this skill requires
4. Body is free-form markdown with numbered steps
5. Use the `skill_create` tool to generate the file automatically
6. Keep skills focused — one task per skill
""",

    'software-dev/rest-graphql-debug.md': """Debug REST and GraphQL APIs systematically.

## REST Debugging
1. Understand the endpoint: method, path, headers, body schema
2. Start with a minimal curl to verify connectivity
3. Check: status codes, response headers, content-type, CORS
4. Common fixes: add auth header, fix content-type, escape JSON properly

## GraphQL Debugging
1. Start with `{ __schema { types { name } } }` to verify endpoint
2. Use `{ __typename }` on any type for validation
3. Check query complexity — batch if hitting limits
4. Handle errors by inspecting `errors[]` array in response

## Pattern
```bash
# REST
curl -X GET "https://api.example.com/v1/resource" -H "Authorization: Bearer $TOKEN"

# GraphQL
curl -X POST "https://api.example.com/graphql" -H "Content-Type: application/json" -d '{"query":"{ users { id name } }"}'
```
""",

    'software-dev/web-development.md': """Apply modern web development patterns:

## Best Practices
1. **Semantic HTML**: Use `<header>`, `<nav>`, `<main>`, `<section>`, `<article>`, `<footer>`
2. **Modern CSS**: Grid, Flexbox, Custom Properties, Container Queries
3. **Responsive**: Mobile-first, `clamp()`, `min()`, `max()` for fluid sizing
4. **Accessibility**: WCAG 2.2 AA — proper contrast (4.5:1), focus indicators, ARIA labels
5. **Performance**: Lazy loading images, code splitting, reduce reflows
6. **Security**: CSP headers, input sanitization, XSS prevention

## Framework Notes (Svelte/SvelteKit)
- Use `$:` reactive statements
- Scoped CSS by default — use `:global()` for global styles
- Store pattern: `writable()`, `derived()`, `get()`
""",

    'software-dev/security.md': """Perform security audit following OWASP Top 10.

## Checklist
1. **Dependencies**: Check for CVEs (`npm audit`, `pip audit`)
2. **Secrets**: Scan for hardcoded API keys, tokens (`grep -r "api_key\|password\|secret"`)
3. **Input Validation**: SQL injection, XSS, command injection vectors
4. **Authentication**: JWT handling, session security, password hashing (bcrypt/argon2)
5. **Authorization**: RBAC, missing permission checks, IDOR
6. **Infrastructure**: TLS 1.3, CORS whitelist, security headers (CSP, HSTS, X-Frame-Options)

## Output Format
```
[SEVERITY] file:line — description — recommendation
```
Severities: CRITICAL | HIGH | MEDIUM | LOW | INFO
""",

    # ─── github ────────────────────────────────────────────────
    'github/github-code-review.md': """Review GitHub Pull Requests with structured feedback.

## Process
1. Fetch PR diff: `gh pr view <number> --json files,body,comments`
2. Review each changed file for:
   - Logic errors and edge cases
   - Code style consistency with project
   - Missing test coverage
   - Security vulnerabilities
3. `github_review_pr` to post review with APPROVE | REQUEST_CHANGES | COMMENT
4. Classify: CRITICAL (blocks merge) | WARNING (should fix) | SUGGESTION (nice-to-have)

## Review Template
```
### Summary
[Overall assessment in 2-3 sentences]

### Findings
- CRITICAL: [file:line] — Description
- WARNING: [file:line] — Description
- SUGGESTION: [file:line] — Description

### Positives
- What's done well
```
""",

    'github/github-auth.md': """Configure GitHub authentication for Nova.

## Steps
1. Check status: `gh auth status`
2. If not authenticated: `gh auth login` (interactive) or set `GITHUB_TOKEN` env
3. Required token scopes: `repo`, `workflow`, `read:org`, `pull_requests`
4. For SSH: check `~/.ssh/id_ed25519.pub`, add via `gh ssh-key add < ~/.ssh/id_ed25519.pub`
5. Verify: `gh api /user` should return user data

## Config
Add to `.env`:
```
GITHUB_TOKEN=ghp_xxxxxxxxxxxxxxxxxxxx
```
Store securely — never expose in logs or output.
""",

    'github/github-pr-workflow.md': """Full PR lifecycle management.

## Create
```bash
git checkout -b feat/description
# make changes, commit
git push -u origin HEAD
gh pr create --title "feat(scope): description" --body "## Summary\n\n- Change 1\n- Change 2" --label enhancement
```

## Manage
- Request reviewers: `gh pr edit <number> --add-reviewer @user`
- Add labels: `gh pr edit <number> --add-label bug,urgent`
- Check CI: `gh pr checks <number>`

## Merge
- Squash merge: `gh pr merge <number> --squash --subject "feat(scope): desc"`
- Delete branch: `git push origin --delete feat/description`
- Use `github_pr_create` tool for automated PR creation
""",

    'github/github-issues.md': """Manage GitHub issues via API.

## Operations
- **List**: `gh issue list --limit 20 --state open` or `github_list_issues`
- **Create**: `github_create_issue` with title, body, labels
- **Update**: `gh issue edit <number> --add-label triage`
- **Close**: `gh issue close <number> --comment "Fixed in #123"`
- **Search**: `gh issue search "label:bug sort:updated-desc"`

## Labels
bug, enhancement, documentation, good-first-issue, help-wanted, question, wontfix, duplicate, triage

Use `github_create_issue` tool for automated issue creation.
""",

    'github/github-repo-management.md': """Manage repository settings and collaborators.

## Create Repo
`gh repo create my-repo --public --clone`

## Settings
```bash
gh repo edit --enable-issues=true --enable-wiki=false
gh repo edit --default-branch main --allow-forking=true
```

## Branch Protection
Require PR reviews, status checks, and linear history:
- Settings → Branches → Add rule for `main`
- Require pull request reviews before merging
- Require status checks to pass

## Secrets & Variables
```bash
gh secret set MY_SECRET --body "value"
gh variable set MY_VAR --body "value"
```

Use `github_search` tool to find repos by query.
""",

    # ─── research ───────────────────────────────────────────────
    'research/arxiv.md': """Search arXiv for academic papers.

## Usage
1. Use `arxiv_search` tool with a search query
2. Results include: title, authors, abstract, URL
3. For PDF access: `https://arxiv.org/pdf/<id>.pdf`

## API Details
Query format: `search_query=all:<term>`
Sort options: relevance, lastUpdatedDate, submittedDate
Max results: 10 (default), up to 100

## Python Parse Snippet
```python
import xml.etree.ElementTree as ET
ns = {'a': 'http://www.w3.org/2005/Atom'}
for entry in ET.fromstring(xml).findall('a:entry', ns):
    title = entry.find('a:title', ns).text.strip()
```
""",

    'research/research-paper-writing.md': """Structure and write academic papers.

## Sections
1. **Title**: Descriptive, 10-20 words
2. **Abstract**: 150-250 words — problem, method, key results, implications
3. **Introduction**: Context → gap → contribution → outline
4. **Related Work**: Position against existing literature with citations
5. **Method**: Reproducible, with algorithm pseudocode
6. **Experiments**: Setup, baselines, datasets, metrics, results tables
7. **Discussion**: Analysis, limitations, ablation studies
8. **Conclusion**: Summary, broader impact, future work

## Tools
- LaTeX: `\cite{key}`, `\ref{fig:1}`, `\section{Introduction}`
- Typst: `#cite(<key>)`, `@fig1`, `== Introduction`
- Generate .bib from DBLP or Semantic Scholar
""",

    'research/llm-wiki.md': """Build a structured wiki from codebase analysis.

## Process
1. Scan the codebase for key modules, APIs, and architecture
2. Use `wiki_build` tool to generate wiki template files
3. Fill in architecture, API reference, config docs, dev guide
4. Cross-link pages with markdown `[links](page.md)`
5. Use `wiki_search` to find content across wiki

## Structure
```
docs/wiki/
  index.md        # Overview + navigation
  architecture.md # System design
  api.md          # API reference
  config.md       # Configuration
  development.md  # Setup, build, deploy
```
""",

    # ─── note-taking ────────────────────────────────────────────
    'note-taking/obsidian.md': """Work with Obsidian vault notes.

## Usage
1. Vault path defaults to `~/obsidian-vault` — configure in settings if different
2. Create notes with `obsidian_create_note` tool
3. Use `[[WikiLinks]]` for internal references
4. Use `#tags` for categorization
5. Frontmatter: tags, created date, aliases

## Best Practices
- One idea per note (Zettelkasten)
- Link generously — build a connected graph
- Use folders for broad categories
- Review and update MOC (Map of Content) files
""",

    # ─── creative ──────────────────────────────────────────────
    'creative/architecture-diagram.md': """Generate dark-theme SVG architecture diagrams.

## Usage
Use `diagram_architecture` tool with:
- title: Diagram title
- components: Array of component names (top to bottom)
- outputFile: Optional custom path

## Style Guide
- Background: `#0a0f10` dark
- Boxes: `#131a1c` with `#2dd4bf` teal borders
- Text: `#ece6da` warm white
- Arrows: Teal dashed lines with arrowheads
- Glow filter on boxes for depth

The output is an HTML file with inline SVG.
""",

    'creative/excalidraw.md': """Create hand-drawn style diagrams using Excalidraw format.

## Usage
Use `excalidraw_generate` tool with:
- title: Diagram name
- elements: Array of {type, x, y, width, height, label}
- outputFile: Optional .excalidraw path

## Element Types
- rectangle: Box for components
- diamond: Decision points
- arrow: Connections (use line elements)
- text: Labels inside boxes

Open the .excalidraw file at https://excalidraw.com or use the VSCode extension.
""",

    # ─── smart-home ────────────────────────────────────────────
    'smart-home/openhue.md': """Control Philips Hue smart lights.

## Setup
1. Find bridge: `curl -s https://discovery.meethue.com/`
2. Get API key: Press bridge button, then `curl -X POST http://<bridge>/api -d '{"devicetype":"nova"}'`
3. Set env: `HUE_BRIDGE_IP` and `HUE_API_KEY`

## Commands
Use `openhue_control` tool:
- `action: "on"` or `"off"` — lightId optional
- `action: "brightness" value: 1-254`
- `action: "color" value: 0-65535`

## Safety
Never change light states without explicit user approval.
""",

    # ─── media ─────────────────────────────────────────────────
    'media/spotify.md': """Search and control Spotify.

## Setup
1. Get API credentials from Spotify Developer Dashboard
2. Set `SPOTIFY_ACCESS_TOKEN` in `.env`
3. For user operations, need OAuth flow with proper scopes

## Operations
- **Search**: `spotify_search` tool — tracks, artists, albums
- Format: track name, artist, album, popularity
- Rate limits: 180 requests per minute

## Scopes
- `user-read-playback-state` — current playback
- `user-modify-playback-state` — play/pause/skip
- `playlist-read-private` — user playlists
""",

    'media/youtube-content.md': """Search and retrieve YouTube content.

## Setup
1. Enable YouTube Data API v3 in Google Cloud Console
2. Get API key and set `YOUTUBE_API_KEY` in `.env`

## Operations
- **Search**: `youtube_search` tool — search by query
- Returns: title, channel, URL, publish date
- Daily quota: 10,000 units (search = 100 units each)

## Transcript
For transcripts, use the video ID with `youtubetranscript.com` API:
`https://youtubetranscript.com/?v=<videoId>`
""",

    'media/gif-search.md': """Search for GIFs to use in chat UI.

## Usage
Use `gif_search` tool with a search query:
- Returns: GIF title, URL, preview URL
- Uses Tenor API — no API key required
- Results are content-filtered
- Works great for chat reactions and visual responses

## Example
"Search for a celebration GIF" → `gif_search(query: "celebration")`
""",

    # ─── social-media ──────────────────────────────────────────
    'social-media/xurl.md': """Interact with X/Twitter API v2.

## Setup
1. Get credentials from Twitter Developer Portal
2. Set `TWITTER_BEARER_TOKEN` in `.env` (read-only)
3. For posting: need OAuth 1.0a user context tokens

## Operations
- **Search**: `twitter_search` — recent tweets by query
- Read-only by default — confirm before any action
- Rate limits: 450 requests per 15 min (app-only)

## Results Format
Tweet text | likes | retweets | timestamp
""",

    # ─── data-science ──────────────────────────────────────────
    'data-science/jupyter-live-kernel.md': """Execute code in Jupyter kernels.

## Operations
- **List kernels**: `jupyter_list_kernels` tool
- Create notebooks: Manually create `.ipynb` JSON files
- Execute: Use shell to run `jupyter nbconvert --execute`

## Notebook Format
```json
{
  "cells": [{
    "cell_type": "code",
    "execution_count": null,
    "source": ["print('hello')"],
    "outputs": []
  }],
  "metadata": {
    "kernelspec": { "display_name": "Python 3", "language": "python", "name": "python3" }
  },
  "nbformat": 4,
  "nbformat_minor": 5
}
```

## Common Kernels
- python3 — Python
- ir — R
- julia-1.x — Julia
""",

    # ─── mcp ───────────────────────────────────────────────────
    'mcp/native-mcp.md': """Integrate MCP (Model Context Protocol) servers.

## About MCP
MCP provides a standardized interface for AI models to interact with external tools and data sources.

## Configuration
Add to config:
```json
{
  "mcpServers": {
    "filesystem": {
      "command": "npx",
      "args": ["-y", "@modelcontextprotocol/server-filesystem", "/path"]
    }
  }
}
```

## In Nova
- Use `/mcp reload` to refresh MCP servers
- MCP tools appear alongside Nova's built-in tools
- Build custom MCP servers for Nova-specific capabilities

## Popular Servers
- `@modelcontextprotocol/server-filesystem` — file operations
- `@modelcontextprotocol/server-github` — GitHub API
- `@modelcontextprotocol/server-brave-search` — web search
""",

    # ─── email ─────────────────────────────────────────────────
    'email/himalaya.md': """Manage email via Himalaya CLI.

## Setup
1. Install: `cargo install himalaya` or `npm -g install himalaya`
2. Configure in `~/.config/himalaya/config.toml`
3. Generate app password for Gmail/Outlook

## Commands
```bash
himalaya list -f INBOX -s 20
himalaya read <id>
himalaya send --to "user@example.com" --subject "Subject" --body "Content"
himalaya search "keyword"
```

## Config Example
```toml
[account.default]
email = "user@example.com"
imap_host = "imap.gmail.com"
imap_port = 993
smtp_host = "smtp.gmail.com"
smtp_port = 465
```
""",

    'email/agentmail.md': """Send and receive emails via AgentMail API.

## Setup
1. Get API key from AgentMail dashboard
2. Set `AGENTMAIL_API_KEY` in `.env`

## Send Email
```bash
curl -X POST https://api.agentmail.io/v1/send \
  -H "Authorization: Bearer $AGENTMAIL_API_KEY" \
  -H "Content-Type: application/json" \
  -d '{"to": ["user@example.com"], "subject": "Hello", "text": "Body"}'
```

Read inbox: `GET https://api.agentmail.io/v1/inbox`

No SMTP/IMAP setup needed — works anywhere with HTTP.
Supports HTML emails, attachments, CC/BCC, reply threading.
""",

    # ─── productivity ──────────────────────────────────────────
    'productivity/google-workspace.md': """Integrate with Google Workspace APIs.

## Setup
1. Enable APIs in Google Cloud Console: Gmail, Docs, Sheets, Drive
2. Create OAuth 2.0 credentials or service account
3. Set `GOOGLE_CLIENT_ID`, `GOOGLE_CLIENT_SECRET`, `GOOGLE_REFRESH_TOKEN`

## Key Endpoints
- Gmail: `GET /gmail/v1/users/me/messages?maxResults=10`
- Sheets: `GET /sheets/v4/spreadsheets/<id>/values/Sheet1!A1:E5`
- Drive: `GET /drive/v3/files?q=mimeType contains 'folder'`
- Docs: `GET /docs/v1/documents/<id>`

## Usage
Read-only by default. Confirm before any write operation.
""",

    'productivity/linear.md': """Interact with Linear project management.

## Setup
1. Get API key from Linear → Settings → API
2. Set `LINEAR_API_KEY` in `.env`
3. API endpoint: `https://api.linear.app/graphql`

## Operations
- List issues: GraphQL query `{ issues { nodes { title state { name } } } }`
- Create issue: `mutation { issueCreate(input: {title: "...", teamId: "..."}) { success } }`
- Update status, priority, assignee

## Format
```bash
curl -s -X POST https://api.linear.app/graphql \
  -H "Authorization: $LINEAR_API_KEY" \
  -H "Content-Type: application/json" \
  -d '{"query":"query { issues { nodes { title } } }"}'
```
""",

    'productivity/airtable.md': """Work with Airtable bases and records.

## Setup
1. Get API key from Airtable Account page
2. Set `AIRTABLE_API_KEY` in `.env`
3. Base ID from URL: `https://airtable.com/<baseId>/...`

## Operations
- List records: `GET /v0/<baseId>/<tableName>?maxRecords=100`
- Create: `POST /v0/<baseId>/<tableName>` with fields
- Update: `PATCH /v0/<baseId>/<tableName>/<recordId>`
- Delete: `DELETE /v0/<baseId>/<tableName>/<recordId>`

API base: `https://api.airtable.com`
Supports filtering by formula, sorting, pagination.
""",

    'productivity/powerpoint.md': """Generate PPTX presentations programmatically.

## Setup
Install python-pptx: `pip install python-pptx`

## Usage
Use `powerpoint_generate` tool:
- title: Presentation title
- slides: Array of {title, content: [bullets]}
- outputFile: Optional custom path

## Features
The generated presentation includes:
- Title slides per section
- Bullet point content
- Consistent styling
- Works on Windows, macOS, Linux

## Direct Python
```python
from pptx import Presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Title"
prs.save("output.pptx")
```
""",

    'productivity/ocr-and-documents.md': """Extract text from images and scanned documents.

## Setup
Install Tesseract OCR:
- Windows: Download from UB-Mannheim GitHub
- macOS: `brew install tesseract`
- Linux: `apt install tesseract-ocr`

## Usage
Use `ocr_extract` tool:
- imagePath: Path to image file
- language: OCR language code (eng, fra, deu, etc.)

## Supported Formats
- PNG, JPG, TIFF, BMP
- Multi-page TIFF
- PDF (convert to images first)

```python
import pytesseract
from PIL import Image
text = pytesseract.image_to_string(Image.open('scan.png'), lang='eng')
```
""",

    # ─── devops ────────────────────────────────────────────────
    'devops/kanban-orchestrator.md': """Orchestrate tasks with Kanban methodology.

## Board Structure
| Backlog | To Do | In Progress | Review | Done |
|---------|-------|-------------|--------|------|

## Workflow
1. Pull from backlog (respect WIP limits)
2. Update status as work progresses
3. Track blockers and dependencies
4. Daily summary with cycle time and throughput

## In Nova
Use `/tasks` to manage the task board:
- `/tasks` — list all tasks
- `/tasks todo` — backlog and todo items
- `/tasks in-progress` — current work
- `/tasks done` — completed items

Track blockers with `/task update <id> --blocked-by <other-id>`.
""",

    # ─── blockchain ────────────────────────────────────────────
    'blockchain/evm.md': """Interact with EVM-compatible blockchains.

## Setup
Set `ETH_RPC_URL` in `.env` (e.g., Alchemy, Infura)

## Operations
- Query balances, read contract state (read-only by default)
- Use ethers.js or direct RPC: `curl -X POST $ETH_RPC_URL -d '{"jsonrpc":"2.0","method":"eth_blockNumber","params":[],"id":1}'`
- For contracts, need ABI — fetch from Etherscan
- Use testnet (Sepolia, Goerli) for testing

## Safety
- Read-only by default
- Never expose private keys
- Confirm before any transaction
- Always test on testnet first
""",

    'blockchain/solana.md': """Interact with Solana blockchain.

## Setup
Set `SOLANA_RPC_URL` in `.env` (default: https://api.mainnet-beta.solana.com)

## Operations
```javascript
const { Connection, PublicKey, LAMPORTS_PER_SOL } = require('@solana/web3.js');
const conn = new Connection(process.env.SOLANA_RPC_URL);
const balance = await conn.getBalance(new PublicKey(address));
```

## Networks
- Mainnet: `https://api.mainnet-beta.solana.com`
- Devnet: `https://api.devnet.solana.com`
- Testnet: `https://api.testnet.solana.com`

## Safety
Read-only by default. Confirm before any transaction.
Use devnet for testing.
""",

    'blockchain/hyperliquid.md': """Interact with Hyperliquid DEX.

## Public API (no auth needed)
```bash
# All mid prices
curl -s -X POST https://api.hyperliquid.xyz/info -d '{"type":"allMids"}'

# Order book
curl -s -X POST https://api.hyperliquid.xyz/info -d '{"type":"l2Book","coin":"BTC"}'

# User state
curl -s -X POST https://api.hyperliquid.xyz/info -d '{"type":"clearinghouseState","user":"0x..."}'
```

## Safety
- Read-only by default (info endpoint)
- Trading requires wallet private key — never store in config
- Never trade without explicit user approval
- Test with small amounts first
""",

    # ─── health ────────────────────────────────────────────────
    'health/health.md': """Track health and fitness metrics.

## Data Format
Track daily in markdown:
```markdown
# Health Log — 2026-05-19
- Steps: 8,234
- Sleep: 7.5h
- Water: 6 glasses
- Workout: Running 30min
- Calories: ~2,100
- Notes: Felt energetic
```

## Features
- Daily logging
- Weekly summaries with averages
- Trend detection
- All data stored locally

## Metrics
Steps, sleep hours, water intake, workouts, nutrition, mood
""",

    # ─── finance ───────────────────────────────────────────────
    'finance/finance.md': """Fetch financial data from public APIs.

## Free APIs (no key required)
- Yahoo Finance: `https://query1.finance.yahoo.com/v8/finance/chart/AAPL`
- CoinGecko (crypto): `https://api.coingecko.com/api/v3/simple/price?ids=bitcoin&vs_currencies=usd`

## Premium APIs (set in .env)
- `ALPHA_VANTAGE_KEY` — stocks, forex
- `POLYGON_KEY` — real-time and historical
- `FINNHUB_KEY` — news, sentiment

## Format
Symbol | Price | Change % | Volume | High/Low | Timestamp

Cache results for 60 seconds to avoid rate limits.
""",

    # ─── sw-dev-opt ────────────────────────────────────────────
    'sw-dev-opt/eslint.md': """Configure and run ESLint for code quality.

## Setup
```bash
npm install --save-dev eslint @eslint/js
```

## Config
```javascript
import js from '@eslint/js';
export default [
  js.configs.recommended,
  { rules: {
    'no-unused-vars': 'warn',
    'no-console': 'warn',
    'prefer-const': 'error',
    'eqeqeq': 'error',
  }}
];
```

## Run
```bash
npx eslint src/ --fix
```

## TypeScript + Svelte
- TypeScript: `typescript-eslint` parser + plugin
- Svelte: `eslint-plugin-svelte`
""",

    'sw-dev-opt/prettier.md': """Format code with Prettier for consistent style.

## Setup
```bash
npm install --save-dev prettier
```

## Config (.prettierrc)
```json
{
  "semi": true,
  "singleQuote": true,
  "tabWidth": 2,
  "trailingComma": "all",
  "printWidth": 100
}
```

## Run
```bash
npx prettier --write .
```

## Plugins
- `prettier-plugin-svelte` — Svelte files
- `prettier-plugin-tailwindcss` — sort Tailwind classes
- `eslint-config-prettier` — avoid conflicts with ESLint
""",

    'sw-dev-opt/typescript.md': """TypeScript best practices and configuration.

## Strict Config
```json
{
  "compilerOptions": {
    "strict": true,
    "noUncheckedIndexedAccess": true,
    "noImplicitReturns": true
  }
}
```

## Best Practices
- Prefer `interface` for object shapes, `type` for unions
- Use `unknown` instead of `any`, narrow with type guards
- Generics with meaningful names: `TItem`, `TKey`, `TValue`
- Utility types: `Partial<T>`, `Pick<T,K>`, `Omit<T,K>`, `Record<K,V>`
- Discriminated unions for state machines
- Branded types for type-safe IDs: `type UserId = string & { __brand: 'UserId' }`
""",

    # ─── community ─────────────────────────────────────────────
    'community/community-download.md': """Download and install community-contributed skills.

## Usage
Use `community_skills_download` tool:
- repo: GitHub repo (default: cheetahclaws/skills)
- skill: Specific skill name to download
- category: Target category folder

## Examples
```bash
# Download all from cheetahclaws/skills
community_skills_download()

# Download a specific skill
community_skills_download(repo: "cheetahclaws/skills", skill: "my-skill", category: "research")

# Download from a custom repo
community_skills_download(repo: "myuser/my-skills")
```

## After Download
Run `skills_list` to see the newly installed skills.
Skills are automatically available to all agents.
""",
}

# Update files
updated = 0
for path, body in bodies.items():
    fp = os.path.join(SKILLS_DIR, path)
    if not os.path.exists(fp):
        print(f"  SKIP (not found): {path}")
        continue
    
    with open(fp, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Find the body separator (after frontmatter)
    if content.startswith('---'):
        end_idx = content.find('---', 3)
        if end_idx != -1:
            frontmatter = content[:end_idx + 3]
            # Write: frontmatter + new body
            with open(fp, 'w', encoding='utf-8') as f:
                f.write(frontmatter + '\n' + body.strip() + '\n')
            updated += 1
            print(f"  OK: {path}")

print(f'\nUpdated {updated} of {len(bodies)} skill files')
